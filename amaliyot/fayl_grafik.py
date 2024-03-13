import os
import docx
from PyPDF2 import PdfReader
from pptx import Presentation
import tkinter as tk
from tkinter import filedialog
root = tk.Tk()
def aniqlash(file_name): 
    try:
        # Faylni hajmini aniqlash qismi
        fayl_hajmi = os.path.getsize(file_name)
        result.insert(tk.END, f"Fayl yo'li: {file_name}")
        if file_name:
            file_size(file_name)
    # FileNotFoundError xatolik chiqsa
    except FileNotFoundError:
        result.insert(tk.END, f"{file_name} fayli topilmadi.\n")

def file_size(file_name):
    fayl_hajmi = os.path.getsize(file_name)
    result1.insert(tk.END, f"Fayl hajmi: {fayl_hajmi} bayt\n")
    # Fayllarni qanday fayl turini bilish qismi
    if file_name.lower().endswith(('.docx', '.doc')):
        doc = docx.Document(file_name)
        qatorlar_soni = sum(len(paragraphs.runs) for paragraphs in doc.paragraphs)
        result1.insert(tk.END, f"Qator: {qatorlar_soni} qator mavjud.\n")

        # Fayllarni qanday fayl turini bilish qismi
    elif file_name.lower().endswith('.pdf'):
         with open(file_name, 'rb') as f:
            pdf_reader = PdfReader(f)
            qatorlar_soni = len(pdf_reader.pages)
            result1.insert(tk.END, f"Qator: {qatorlar_soni} qator mavjud.\n")

        # Fayllarni qanday fayl turini bilish qismi
    elif file_name.lower().endswith('.pptx'):
        prs = Presentation(file_name)
        qatorlar_soni = sum(len(slide.shapes) for slide in prs.slides)
        result1.insert(tk.END, f"Qator: {qatorlar_soni} qator mavjud.\n")

        # Fayl mos kelmasa
    else:
        result1.insert(tk.END, f"Qator: {qatorlar_soni} qator mavjud.\n")

def file_choose():
    file_path = filedialog.askopenfilename()
    if file_path:
        aniqlash(file_path)
    elif file_path:
        file_size(file_path)

# Tkinter main window

root.title("Fayl ma'lumotlarini ko'rish dasturi")
browse_button = tk.Button(root, text="Faylni ko'rib chiqish", command=file_choose, bg='green', fg='white')
browse_button.pack(pady=10)

paned_window = tk.PanedWindow(root, orient=tk.HORIZONTAL)
paned_window.pack(fill=tk.BOTH, expand=True)

result = tk.Text(paned_window, height=15, width=40, bg='blue', fg='white')
result.pack(padx=10)
paned_window.add(result)

result1 = tk.Text(paned_window, height=15, width=40, bg='yellow', fg='black')
result1.pack(padx=10)
paned_window.add(result1)

button_frame = tk.Frame(root)
button_frame.pack(side=tk.RIGHT, fill=tk.X)

btn_quit = tk.Button(button_frame, text='Chiqish', command=root.destroy, bg='red', fg='white')
btn_quit.pack(side=tk.LEFT, padx=10, pady=5)

btn_scan = tk.Button(button_frame, text='Skaner qilish', command=file_choose, bg='blue', fg='white')
btn_scan.pack(side=tk.LEFT, padx=10, pady=5)

root.mainloop()
