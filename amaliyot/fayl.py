import os
import docx 
from PyPDF2 import PdfReader 
from pptx import Presentation

def aniqlash(fayl_nom):
    try:
        # faylni hajmini aniqlash qismi
        fayl_hajmi = os.path.getsize(fayl_nom)
        print(f"{fayl_nom} hajmi: {fayl_hajmi} bayt")

        # fayllarni qanday fayl turini bilish qismi
        if fayl_nom.lower().endswith(('.docx', '.doc')):
            doc = docx.Document(fayl_nom)
            qatorlar_soni = sum(len(paragraphs.runs) for paragraphs in doc.paragraphs)
            print(f"{fayl_nom}da {qatorlar_soni} qator mavjud.")
        
        # fayllarni qanday fayl turini bilish qismi
        elif fayl_nom.lower().endswith('.pdf'):
            with open(fayl_nom, 'rb') as f:
                pdf_reader = PdfReader(f)
                qatorlar_soni = len(pdf_reader.pages)
                print(f"{fayl_nom}da {qatorlar_soni} qator mavjud.")
        
        # fayllarni qanday fayl turini bilish qismi
        elif fayl_nom.lower().endswith('.pptx'):
            prs = Presentation(fayl_nom)
            qatorlar_soni = sum(len(slide.shapes) for slide in prs.slides)
            print(f"{fayl_nom}da {qatorlar_soni} qator mavjud.")
        
        # fayl mos kelmasa
        else:
            print(f"{fayl_nom} Ushbu fayl formati qo'llab-quvvatlanmaydi.")

    # FileNotFoundError xatolik chiqsa
    except FileNotFoundError:
        print(f"{fayl_nom} fayli topilmadi.")


if __name__ == "__main__":
    files = os.listdir()
    #bizda mavjud fayl
    for fayl_nom in files:
        fayl_nom = os.path.join(fayl_nom)
        aniqlash(fayl_nom)
