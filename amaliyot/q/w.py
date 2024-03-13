import tkinter as tk
from tkinter import filedialog, messagebox
import os
import docx 
from PyPDF2 import PdfReader 
from pptx import Presentation

def fayllar_royxati_ochish(fayllar):
    royxat = "\n".join(fayllar)
    messagebox.showinfo("Papka Fayllari", royxat)

def aniqlash(fayl_nom):
    try:
        # Faylni hajmini aniqlash
        fayl_hajmi = os.path.getsize(fayl_nom)
        # Faylning turini aniqlash
        if fayl_nom.lower().endswith(('.docx', '.doc')):
            doc = docx.Document(fayl_nom)
            qatorlar_soni = sum(len(paragraphs.runs) for paragraphs in doc.paragraphs)
            zararli = False
        elif fayl_nom.lower().endswith('.pdf'):
            with open(fayl_nom, 'rb') as f:
                pdf_reader = PdfReader(f)
                qatorlar_soni = len(pdf_reader.pages)
                zararli = False
        elif fayl_nom.lower().endswith('.pptx'):
            prs = Presentation(fayl_nom)
            qatorlar_soni = sum(len(slide.shapes) for slide in prs.slides)
            zararli = False
        elif fayl_nom.lower().endswith('.bat'):
            with open(fayl_nom, 'r') as f:
                qatorlar_soni = sum(1 for line in f)
                zararli = True
        else:
            print(f"{fayl_nom} Ushbu fayl formati qo'llab-quvvatlanmaydi.")
            qatorlar_soni = 0
            zararli = False
        return fayl_hajmi, qatorlar_soni, zararli

    except FileNotFoundError:
        print(f"{fayl_nom} fayli topilmadi.")
        return 0, 0, False

# def aniqlash(fayl_nom):
#     try:
#         # Faylni hajmini aniqlash
#         fayl_hajmi = os.path.getsize(fayl_nom)
#         # Faylning turini aniqlash
#         if fayl_nom.lower().endswith(('.docx', '.doc')):
#             doc = docx.Document(fayl_nom)
#             qatorlar_soni = sum(len(paragraphs.runs) for paragraphs in doc.paragraphs)
#             zararli = False
#         elif fayl_nom.lower().endswith('.pdf'):
#             with open(fayl_nom, 'rb') as f:
#                 pdf_reader = PdfReader(f)
#                 qatorlar_soni = len(pdf_reader.pages)
#                 zararli = False
#         elif fayl_nom.lower().endswith('.pptx'):
#             prs = Presentation(fayl_nom)
#             qatorlar_soni = sum(len(slide.shapes) for slide in prs.slides)
#             zararli = False
#         else:
#             print(f"{fayl_nom} Ushbu fayl formati qo'llab-quvvatlanmaydi.")
#             qatorlar_soni = 0
#             zararli = False
#         return fayl_hajmi, qatorlar_soni, zararli

#     except FileNotFoundError:
#         print(f"{fayl_nom} fayli topilmadi.")
#         return 0, 0, False

def skaner():
    global selected_folder
    if selected_folder:
        # Oynani chiqarish uchun ikkinchi oyna
        skaner_window = tk.Toplevel(root)
        skaner_window.title("Skaner Natijalari")

        # Canvas yaratish
        canvas = tk.Canvas(skaner_window)
        canvas.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)

        # Skrollbar yaratish
        scrollbar = tk.Scrollbar(skaner_window, orient=tk.VERTICAL, command=canvas.yview)
        scrollbar.pack(side=tk.RIGHT, fill=tk.Y)

        # Canvas va skrollbar o'rtasida alohida bo'ylab bog'liq
        canvas.configure(yscrollcommand=scrollbar.set)

        # Fayllar ro'yxatini joylash uchun ramka
        frame = tk.Frame(canvas)
        canvas.create_window((0, 0), window=frame, anchor=tk.NW)

        for fayl_nom in os.listdir(selected_folder):
            fayl_nom = os.path.join(selected_folder, fayl_nom)
            hajmi, qatorlar_soni, zararli = aniqlash(fayl_nom)

            # Natijalarni chiqarish
            natija_str = f"Fayl nomi: {os.path.basename(fayl_nom)}\n"
            natija_str += f"Hajmi: {hajmi} bayt\n"
            natija_str += f"Qatorlar soni: {qatorlar_soni}\n"
            natija_str += "Xavfli" if zararli else "Xavfsiz"
            natija_label = tk.Label(frame, text=natija_str)
            natija_label.pack()

        # Fayllarni chiqarish tugmasi
        chiqarish_btn = tk.Button(skaner_window, text="Chiqarish", command=skaner_window.destroy)
        chiqarish_btn.pack()

        messagebox.showinfo("Tugma bosilganda", "Fayllar skanlandi va natijalar chiqarildi.")
    else:
        messagebox.showwarning("Diqqat", "Iltimos, avval 'Papka' tugmasini bosing va papkani tanlang.")

# def skaner():
#     global selected_folder
#     if selected_folder:
#         # Oynani chiqarish uchun ikkinchi oyna
#         skaner_window = tk.Toplevel(root)
#         skaner_window.title("Skaner Natijalari")

#         # Canvas va Scrollbar qo'shish
#         natija_frame = tk.Frame(skaner_window)
#         natija_frame.pack(fill="both", expand=True)

#         canvas = tk.Canvas(natija_frame)
#         canvas.pack(side="left", fill="both", expand=True)

#         scrollbar = tk.Scrollbar(natija_frame, orient="vertical", command=canvas.yview)
#         scrollbar.pack(side="right", fill="y")

#         canvas.configure(yscrollcommand=scrollbar.set)
#         canvas.bind("<Configure>", lambda e: canvas.configure(scrollregion=canvas.bbox("all")))

#         # Natijalar ro'yxati uchun yangi ramka
#         natija_list_frame = tk.Frame(canvas)
#         canvas.create_window((0, 0), window=natija_list_frame, anchor="nw")

#         for fayl_nom in os.listdir(selected_folder):
#             fayl_nom = os.path.join(selected_folder, fayl_nom)
#             hajmi, qatorlar_soni, zararli = aniqlash(fayl_nom)

#             # Natijalarni chiqarish
#             natija_str = f"Fayl nomi: {os.path.basename(fayl_nom)}\n"
#             natija_str += f"Hajmi: {hajmi} bayt\n"
#             natija_str += f"Qatorlar soni: {qatorlar_soni}\n"
#             natija_str += "Xavfli" if zararli else "Xavfsiz"
#             natija_label = tk.Label(natija_list_frame, text=natija_str)
#             natija_label.pack()

#         messagebox.showinfo("Tugma bosilganda", "Fayllar skanlandi va natijalar chiqarildi.")
#     else:
#         messagebox.showwarning("Diqqat", "Iltimos, avval 'Papka' tugmasini bosing va papkani tanlang.")

def papka_tanlash():
    global selected_folder
    selected_folder = filedialog.askdirectory()
    if selected_folder:
        fayllar = os.listdir(selected_folder)
        fayllar_royxati_ochish(fayllar)

root = tk.Tk()
root.title("Fayl Skaner")

selected_folder = None

frame = tk.Frame(root)
frame.pack(padx=20, pady=20)

btn_papka = tk.Button(frame, text="Papka", command=papka_tanlash)
btn_papka.pack(side=tk.LEFT, padx=10)

btn_skaner = tk.Button(frame, text="Skaner", command=skaner)
btn_skaner.pack(side=tk.LEFT, padx=10)

btn_ochirish = tk.Button(frame, text="Chiqish", command=root.destroy)
btn_ochirish.pack(side=tk.LEFT, padx=10)

root.geometry("300x150")
root.configure(bg="lightgray")
frame.configure(bg="lightgray")
btn_papka.configure(bg="white")
btn_skaner.configure(bg="white")
btn_ochirish.configure(bg="white")

root.mainloop()
