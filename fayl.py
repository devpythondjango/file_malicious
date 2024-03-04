import os
import docx 
from PyPDF2 import PdfReader 
from pptx import Presentation

def aniqlash(fayl_nom):
    try:
        # faylni hajmini aniqlash qismi
        fayl_hajmi = os.path.getsize(fayl_nom)
        print(f"{fayl_nom} hajmi: {fayl_hajmi} bayt")

        # fayllarni qanday fayl turini bilish qismi
        if fayl_nom.lower().endswith(('.docx', '.doc')):
            doc = docx.Document(fayl_nom)
            qatorlar_soni = sum(len(paragraphs.runs) for paragraphs in doc.paragraphs)
            print(f"{fayl_nom}da {qatorlar_soni} qator mavjud.")
        
        # fayllarni qanday fayl turini bilish qismi
        elif fayl_nom.lower().endswith('.pdf'):
            with open(fayl_nom, 'rb') as f:
                pdf_reader = PdfReader(f)
                qatorlar_soni = len(pdf_reader.pages)
                print(f"{fayl_nom}da {qatorlar_soni} qator mavjud.")
        
        # fayllarni qanday fayl turini bilish qismi
        elif fayl_nom.lower().endswith('.pptx'):
            prs = Presentation(fayl_nom)
            qatorlar_soni = sum(len(slide.shapes) for slide in prs.slides)
            print(f"{fayl_nom}da {qatorlar_soni} qator mavjud.")
        
        # fayl mos kelmasa
        else:
            print(f"{fayl_nom} Ushbu fayl formati qo'llab-quvvatlanmaydi.")

    # FileNotFoundError xatolik chiqsa
    except FileNotFoundError:
        print(f"{fayl_nom} fayli topilmadi.")


if __name__ == "__main__":
    files = os.listdir()
    # #bizda mavjud fayl
    for i in files:
        fayl_nom = i
        aniqlash(fayl_nom)
        
    
# import fitz  # PyMuPDF kutubxonasini o'rnating

# def read_pdf(file_path):
#     with fitz.open(file_path) as pdf_document:
#         text = ""
#         for page_num in range(pdf_document.page_count):
#             page = pdf_document[page_num]
#             text += page.get_text()
#     return text

# pdf_content = read_pdf("fayl.pdf")
# print(pdf_content)

# from docx import Document

# def read_docx(file_path):
#     doc = Document(file_path)
#     text = ""
#     for paragraph in doc.paragraphs:
#         text += paragraph.text + "\n"
#     return text

# docx_content = read_docx("fayl.docx")
# print(docx_content)


# from pptx import Presentation

# def read_pptx(file_path):
#     prs = Presentation(file_path)
#     text = ""
#     for slide in prs.slides:
#         for shape in slide.shapes:
#             if hasattr(shape, "text"):
#                 text += shape.text + "\n"
#     return text

# pptx_content = read_pptx("fayl.pptx")
# print(pptx_content)
