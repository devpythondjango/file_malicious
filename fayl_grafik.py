import os
import docx
from PyPDF2 import PdfReader
from pptx import Presentation
import tkinter as tk
from tkinter import filedialog

def aniqlash(fayl_nom):
    try:
        # Faylni hajmini aniqlash qismi
        fayl_hajmi = os.path.getsize(fayl_nom)
        result_text.insert(tk.END, f"{fayl_nom} hajmi: {fayl_hajmi} bayt\n")

        # Fayllarni qanday fayl turini bilish qismi
        if fayl_nom.lower().endswith(('.docx', '.doc')):
            doc = docx.Document(fayl_nom)
            qatorlar_soni = sum(len(paragraphs.runs) for paragraphs in doc.paragraphs)
            result_text.insert(tk.END, f"{fayl_nom}da {qatorlar_soni} qator mavjud.\n")

        # Fayllarni qanday fayl turini bilish qismi
        elif fayl_nom.lower().endswith('.pdf'):
            with open(fayl_nom, 'rb') as f:
                pdf_reader = PdfReader(f)
                qatorlar_soni = len(pdf_reader.pages)
                result_text.insert(tk.END, f"{fayl_nom}da {qatorlar_soni} qator mavjud.\n")

        # Fayllarni qanday fayl turini bilish qismi
        elif fayl_nom.lower().endswith('.pptx'):
            prs = Presentation(fayl_nom)
            qatorlar_soni = sum(len(slide.shapes) for slide in prs.slides)
            result_text.insert(tk.END, f"{fayl_nom}da {qatorlar_soni} qator mavjud.\n")

        # Fayl mos kelmasa
        else:
            result_text.insert(tk.END, f"{fayl_nom} Ushbu fayl formati qo'llab-quvvatlanmaydi.\n")

    # FileNotFoundError xatolik chiqsa
    except FileNotFoundError:
        result_text.insert(tk.END, f"{fayl_nom} fayli topilmadi.\n")

def browse_file():
    file_path = filedialog.askopenfilename()
    if file_path:
        result_text.delete(1.0, tk.END)  # Clear previous results
        aniqlash(file_path)

# Tkinter main window
root = tk.Tk()
root.title("Fayl ma'lumotlarini ko'rish dasturi")

# Create and pack widgets
browse_button = tk.Button(root, text="Faylni ko'rib chiqish", command=browse_file)
browse_button.pack(pady=10)

result_text = tk.Text(root, height=10, width=50)
result_text.pack()

root.mainloop()
