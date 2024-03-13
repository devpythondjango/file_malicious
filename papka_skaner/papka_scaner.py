import tkinter as tk
from tkinter import filedialog, messagebox
import os
import docx 
from PyPDF2 import PdfReader 
from pptx import Presentation

def fayllar_royxati_ochish(fayllar):
    papka_window = tk.Toplevel(root)
    papka_window.title("Papka Fayllari")

    # Scrollbar qo'shish
    scrollbar = tk.Scrollbar(papka_window)
    scrollbar.pack(side=tk.RIGHT, fill=tk.Y)

    # Text uchun
    fayllar_text = tk.Text(papka_window, yscrollcommand=scrollbar.set, wrap=tk.NONE, width=60, height=20)
    fayllar_text.pack(expand=True, fill=tk.BOTH)

    # Scrollbar va Textni bog'lash
    scrollbar.config(command=fayllar_text.yview)

    fayllar_text.insert(tk.END, "\n".join(fayllar))

def aniqlash(fayl_nom):
    try:
        # Faylni hajmini aniqlash
        fayl_hajmi = os.path.getsize(fayl_nom)
        
        # Faylning turini aniqlash
        if fayl_nom.lower().endswith(('.docx', '.doc')):
            doc = docx.Document(fayl_nom)
            qatorlar_soni = sum(len(paragraphs.runs) for paragraphs in doc.paragraphs)
            zararli = False
        elif fayl_nom.lower().endswith('.pdf'):
            with open(fayl_nom, 'rb') as f:
                pdf_reader = PdfReader(f)
                qatorlar_soni = len(pdf_reader.pages)
                zararli = False
        elif fayl_nom.lower().endswith('.pptx'):
            prs = Presentation(fayl_nom)
            qatorlar_soni = sum(len(slide.shapes) for slide in prs.slides)
            zararli = False
        elif fayl_nom.lower().endswith(('.bat', '.ps1', '.vbs')):
            zararli = True
            # Fayl turi zararli bo'lsa, uni aniqlash uchun qo'shimcha ishlarni bajarish mumkin
            qatorlar_soni = "N/A"
            # Zararli deb belgilaymiz
            messagebox.showinfo("Ogohlantirish", f"{fayl_nom} virus topildi.")
        else:
            # Ushbu fayl formati qo'llab-quvvatlanmaydi
            messagebox.showinfo("Ogohlantirish", f"{fayl_nom} Ushbu fayl formati qo'llab-quvvatlanmaydi.")
            qatorlar_soni = 0
            zararli = False
        
        return fayl_hajmi, qatorlar_soni, zararli

    except FileNotFoundError:
        # Fayl topilmadi
        messagebox.showinfo("Xatolik", f"{fayl_nom} fayli topilmadi.")
        return 0, 0, False

def ochirish(fayl_nom):
    try:
        # Faylni ochirish
        os.remove(fayl_nom)
        messagebox.showinfo("Xabar", f"{fayl_nom} fayli muvaffaqiyatli o'chirildi.")
    except Exception as e:
        # Ochirishda xatolik bo'lgan holat
        messagebox.showerror("Xatolik", f"{fayl_nom} faylini o'chirishda xatolik yuz berdi: {e}")

def skaner():
    global selected_folder
    if selected_folder:
        # Oynani chiqarish uchun ikkinchi oyna
        skaner_window = tk.Toplevel(root)
        skaner_window.title("Skaner Natijalari")

        # Scrollbar qo'shish
        scrollbar = tk.Scrollbar(skaner_window)
        scrollbar.pack(side=tk.RIGHT, fill=tk.Y)

        # Natijalarni chiqarish uchun Text maydonchasi
        natijalar_text = tk.Text(skaner_window, yscrollcommand=scrollbar.set)
        natijalar_text.pack(expand=True, fill=tk.BOTH)

        # Scrollbar va Textni bog'lash
        scrollbar.config(command=natijalar_text.yview)

        zararli_fayllar = []  # Zararli fayllarni saqlash uchun ro'yxat

        for fayl_nom in os.listdir(selected_folder):
            fayl_nom = os.path.join(selected_folder, fayl_nom)
            hajmi, qatorlar_soni, zararli = aniqlash(fayl_nom)

            # Natijalarni chiqarish
            natija_str = f"Fayl nomi: {os.path.basename(fayl_nom)}\n"
            natija_str += f"Hajmi: {hajmi} bayt\n"
            natija_str += f"Qatorlar soni: {qatorlar_soni}\n"
            natija_str += "Xavfli" if zararli else "Xavfsiz"
            natija_str += "\n\n"
            natijalar_text.insert(tk.END, natija_str)

            # Fayl zararli bo'lsa, ochirish tugmasini qo'shamiz va ro'yxatga qo'shamiz
            if zararli:
                ochirish_btn = tk.Button(skaner_window, text=f"{os.path.basename(fayl_nom)} o'chirish", command=lambda f=fayl_nom: ochirish(f))
                ochirish_btn.pack()
                zararli_fayllar.append(fayl_nom)

        # Fayllarni chiqarish tugmasi
        chiqarish_btn = tk.Button(skaner_window, text="Chiqarish", command=skaner_window.destroy)
        chiqarish_btn.pack()

        # Zararli fayllarni ochirish tugmasi
        if zararli_fayllar:
            def ochirish_zararli_fayllar():
                for fayl_nom in zararli_fayllar:
                    os.remove(fayl_nom)
                messagebox.showinfo("Xabar", "Zararli fayllar muvaffaqiyatli o'chirildi.")
                skaner_window.destroy()

            ochirish_btn = tk.Button(skaner_window, text="zararli fayllarni o'chirish", command=ochirish_zararli_fayllar)
            ochirish_btn.pack()

        messagebox.showinfo("Tugma bosilganda", "Fayllar skanlandi va natijalar chiqarildi.")
    else:
        messagebox.showwarning("Diqqat", "Iltimos, avval 'Papka' tugmasini bosing va papkani tanlang.")

# def skaner():
#     global selected_folder
#     if selected_folder:
#         # Oynani chiqarish uchun ikkinchi oyna
#         skaner_window = tk.Toplevel(root)
#         skaner_window.title("Skaner Natijalari")

#         # Scrollbar qo'shish
#         scrollbar = tk.Scrollbar(skaner_window)
#         scrollbar.pack(side=tk.RIGHT, fill=tk.Y)

#         # Natijalarni chiqarish uchun Text maydonchasi
#         natijalar_text = tk.Text(skaner_window, yscrollcommand=scrollbar.set)
#         natijalar_text.pack(expand=True, fill=tk.BOTH)

#         # Scrollbar va Textni bog'lash
#         scrollbar.config(command=natijalar_text.yview)

#         for fayl_nom in os.listdir(selected_folder):
#             fayl_nom = os.path.join(selected_folder, fayl_nom)
#             hajmi, qatorlar_soni, zararli = aniqlash(fayl_nom)

#             # Natijalarni chiqarish
#             natija_str = f"Fayl nomi: {os.path.basename(fayl_nom)}\n"
#             natija_str += f"Hajmi: {hajmi} bayt\n"
#             natija_str += f"Qatorlar soni: {qatorlar_soni}\n"
#             natija_str += "Xavfli" if zararli else "Xavfsiz"
#             natija_str += "\n\n"
#             natijalar_text.insert(tk.END, natija_str)

#             # Fayl zararli bo'lsa, ochirish tugmasini qo'shamiz
#             if zararli:
#                 ochirish_btn = tk.Button(skaner_window, text=f"{os.path.basename(fayl_nom)} o'chirish", command=lambda f=fayl_nom: ochirish(f))
#                 ochirish_btn.pack()

#         # Fayllarni chiqarish tugmasi
#         chiqarish_btn = tk.Button(skaner_window, text="Chiqarish", command=skaner_window.destroy)
#         chiqarish_btn.pack()
        
#         messagebox.showinfo("Tugma bosilganda", "Fayllar skanerlandi va natijalar chiqarildi.")
#     else:
#         messagebox.showwarning("Diqqat", "Iltimos, avval 'Papka' tugmasini bosing va papkani tanlang.")

# def skaner():
#     global selected_folder
#     if selected_folder:
#         # Oynani chiqarish uchun ikkinchi oyna
#         skaner_window = tk.Toplevel(root)
#         skaner_window.title("Skaner Natijalari")

#         # Scrollbar qo'shish
#         scrollbar = tk.Scrollbar(skaner_window)
#         scrollbar.pack(side=tk.RIGHT, fill=tk.Y)

#         # Natijalarni chiqarish uchun Text maydonchasi
#         natijalar_text = tk.Text(skaner_window, yscrollcommand=scrollbar.set)
#         natijalar_text.pack(expand=True, fill=tk.BOTH)

#         # Scrollbar va Textni bog'lash
#         scrollbar.config(command=natijalar_text.yview)

#         for fayl_nom in os.listdir(selected_folder):
#             fayl_nom = os.path.join(selected_folder, fayl_nom)
#             hajmi, qatorlar_soni, zararli = aniqlash(fayl_nom)

#             # Natijalarni chiqarish
#             natija_str = f"Fayl nomi: {os.path.basename(fayl_nom)}\n"
#             natija_str += f"Hajmi: {hajmi} bayt\n"
#             natija_str += f"Qatorlar soni: {qatorlar_soni}\n"
#             natija_str += "Xavfli" if zararli else "Xavfsiz"
#             natija_str += "\n\n"
#             natijalar_text.insert(tk.END, natija_str)

#         # Fayllarni chiqarish tugmasi
#         chiqarish_btn = tk.Button(skaner_window, text="Chiqarish", command=skaner_window.destroy)
#         chiqarish_btn.pack()
        
#         messagebox.showinfo("Tugma bosilganda", "Fayllar skanlandi va natijalar chiqarildi.")
#     else:
#         messagebox.showwarning("Diqqat", "Iltimos, avval 'Papka' tugmasini bosing va papkani tanlang.")

def papka_tanlash():
    global selected_folder
    selected_folder = filedialog.askdirectory()
    if selected_folder:
        fayllar = os.listdir(selected_folder)
        fayllar_royxati_ochish(fayllar)

root = tk.Tk()
root.title("Fayl Skaner")

selected_folder = None

frame = tk.Frame(root)
frame.pack(padx=20, pady=20)

btn_papka = tk.Button(frame, text="Papka", command=papka_tanlash)
btn_papka.pack(side=tk.LEFT, padx=10)

btn_skaner = tk.Button(frame, text="Skaner", command=skaner)
btn_skaner.pack(side=tk.LEFT, padx=10)

btn_ochirish = tk.Button(frame, text="Chiqish", command=root.destroy)
btn_ochirish.pack(side=tk.LEFT, padx=10)

root.geometry("300x150")
root.configure(bg="lightgray")
frame.configure(bg="lightgray")
btn_papka.configure(bg="white")
btn_skaner.configure(bg="white")
btn_ochirish.configure(bg="white")

root.mainloop()
